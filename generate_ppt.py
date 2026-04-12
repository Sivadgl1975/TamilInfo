from pptx import Presentation
import re

# Read the Markdown file
with open('AI_Context_Engineering_Master_Presentation.md', 'r') as f:
    content = f.read()

# Create a new presentation
prs = Presentation()

# Add title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "AI & Context Engineering Master Presentation"
subtitle = title_slide.placeholders[1]
subtitle.text = "A Complete Visual Teaching Guide\nBased on comprehensive AI and context engineering principles\nPresented by [Your Name]\nDate: April 12, 2026"

# Find all slide sections starting with #### Slide
slide_pattern = r'(#### Slide \d+.*?(?=\n#### Slide|\n## |\Z))'
slides = re.findall(slide_pattern, content, re.DOTALL)

# Add each slide
for slide_content in slides:
    lines = slide_content.strip().split('\n')
    title_line = lines[0].replace('#### ', '').strip()
    bullets = [line.strip('- ').strip() for line in lines[1:] if line.strip().startswith('-')]
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_line
    
    if bullets:
        body = slide.placeholders[1]
        tf = body.text_frame
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

# Add PPT Structure Overview
table_match = re.search(r'## 📊 PPT Structure Overview\n(.*?)\n## 🎭', content, re.DOTALL)
if table_match:
    table_text = table_match.group(1).strip()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PPT Structure Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = table_text

# Add Character Design System
char_match = re.search(r'## 🎭 Character Design System\n(.*?)\n## 📑', content, re.DOTALL)
if char_match:
    char_text = char_match.group(1).strip()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Character Design System"
    body = slide.placeholders[1]
    tf = body.text_frame
    for line in char_text.split('\n'):
        if line.strip().startswith('-'):
            p = tf.add_paragraph()
            p.text = line.strip('- ').strip()

# Add Design Specifications
design_match = re.search(r'## 🎨 DESIGN SPECIFICATIONS\n(.*?)\n## 📥', content, re.DOTALL)
if design_match:
    design_text = design_match.group(1).strip()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Design Specifications"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = design_text

# Add Deliverable Checklist
check_match = re.search(r'## 📥 DELIVERABLE CHECKLIST\n(.*?)\Z', content, re.DOTALL)
if check_match:
    check_text = check_match.group(1).strip()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deliverable Checklist"
    body = slide.placeholders[1]
    tf = body.text_frame
    for line in check_text.split('\n'):
        if line.strip().startswith('-'):
            p = tf.add_paragraph()
            p.text = line.strip('- ').strip()

# Save the presentation
prs.save('AI_Context_Engineering_Master_Presentation.pptx')
print("Presentation generated successfully: AI_Context_Engineering_Master_Presentation.pptx")